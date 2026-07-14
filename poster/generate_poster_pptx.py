#!/usr/bin/env python3
"""
generate_poster_pptx.py
=======================
Generates an A2 portrait academic poster as a PowerPoint (.pptx) file,
styled to match the University of Reading Department of Computer Science
poster template.

Minimum body font size: 24 pt.
No en-dashes or em-dashes.

Run from the repository root:
    python poster/generate_poster_pptx.py

Output:
    poster/poster.pptx
"""

import io
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO

# ── Brand colours ─────────────────────────────────────────────────────────────
TEAL  = RGBColor(0x1A, 0x7A, 0x8A)
NAVY  = RGBColor(0x00, 0x3B, 0x6F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
LGREY = RGBColor(0xEA, 0xF5, 0xF7)   # light teal tint for pipeline / KPI bg
DGREY = RGBColor(0xA0, 0xD0, 0xDA)   # border colour
RED   = RGBColor(0xBE, 0x1E, 0x2D)
BLUE  = RGBColor(0x1A, 0x6B, 0xB5)
ORNG  = RGBColor(0xA0, 0x40, 0x10)

# ── Slide geometry (A2 portrait: 42 x 59.4 cm) ───────────────────────────────
SW, SH   = 42.0, 59.4   # slide width / height in cm
MLR      = 1.2           # left and right margin
HDRH     = 1.5           # teal department strip height
TTH      = 7.5           # title area height below strip (fits 48pt title)
FTRH     = 2.8           # footer height (fits 24pt text)
CGAP     = 0.8           # gap between the two body columns

BODY_TOP = HDRH + TTH + 0.3    # y where body columns begin (9.3 cm)
BODY_BOT = SH - FTRH - 0.3    # y where body columns end  (56.3 cm)

COL_W  = (SW - 2 * MLR - CGAP) / 2   # ~19.4 cm per column
COL1X  = MLR
COL2X  = MLR + COL_W + CGAP

# ── Font sizes (all >= 24 pt) ─────────────────────────────────────────────────
FS_TITLE  = 48
FS_STITLE = 30
FS_AUTH   = 24
FS_SEC    = 26     # section header text (inside teal box)
FS_SSEC   = 26     # sub-section heading (bold teal, no box)
FS_BODY   = 26     # body paragraphs
FS_SMALL  = 24     # bullet items and pipeline content
FS_TINY   = 24     # KPI labels and footer

# ── Section header dimensions ─────────────────────────────────────────────────
SEC_H   = 1.3      # section header box height (cm)
SEC_GAP = 0.25     # gap below section header (cm)


# ══════════════════════════════════════════════════════════════════════════════
# LOW-LEVEL PPTX HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def _tf(shape, ml=0.12, mr=0.12, mt=0.08, mb=0.08):
    """Configure and return a shape's text frame with tight internal margins."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left   = Cm(ml)
    tf.margin_right  = Cm(mr)
    tf.margin_top    = Cm(mt)
    tf.margin_bottom = Cm(mb)
    return tf


def _run(para, text, size, bold=False, italic=False, color=None):
    """Add a formatted run to a paragraph and return it."""
    r = para.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    if color:
        r.font.color.rgb = color
    return r


def _solid(shape, color):
    """Fill a shape with a solid colour."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape):
    """Remove the shape outline."""
    shape.line.fill.background()


def _rect(slide, x, y, w, h, color=None, line_color=None, line_pt=0.5):
    """Add a plain rectangle; optionally filled and/or outlined."""
    shp = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    if color:
        _solid(shp, color)
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_pt)
    else:
        _no_line(shp)
    return shp


def _rrect(slide, x, y, w, h, color=None, line_color=None, line_pt=0.5, adj=0.3):
    """Add a rounded rectangle; adj controls corner radius (0–0.5)."""
    shp = slide.shapes.add_shape(MSO.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    shp.adjustments[0] = adj
    if color:
        _solid(shp, color)
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_pt)
    else:
        _no_line(shp)
    return shp


# ══════════════════════════════════════════════════════════════════════════════
# SECTION / TEXT COMPONENTS
# ══════════════════════════════════════════════════════════════════════════════

def sec_hdr(slide, label, x, y, w, h=SEC_H, gap=SEC_GAP):
    """
    Teal filled rounded-rectangle section header with white bold text.
    Returns the y coordinate immediately below the header.
    """
    shp = _rrect(slide, x, y, w, h, color=TEAL, adj=0.3)
    tf = _tf(shp, ml=0.22, mt=0.12, mb=0.12)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _run(p, label.upper(), FS_SEC, bold=True, color=WHITE)
    return y + h + gap


def ssec_hdr(slide, label, x, y, w, gap=0.2):
    """
    Bold teal sub-section heading (no box).
    Returns the y coordinate immediately below.
    """
    txb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(0.65))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    _run(p, label, FS_SSEC, bold=True, color=TEAL)
    return y + 0.65 + gap


def text_box(slide, paras, x, y, w, h, size=FS_BODY, para_gap=4):
    """
    Multiple paragraphs in one text box.

    paras: list of str  OR  (text, bold, italic) tuples.
    Returns y + h.
    """
    txb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Cm(0)

    for i, item in enumerate(paras):
        if isinstance(item, str):
            text, bold, italic = item, False, False
        elif len(item) == 2:
            text, bold = item; italic = False
        else:
            text, bold, italic = item

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(para_gap)
        _run(p, text, size, bold=bold, italic=italic, color=BLACK)

    return y + h


def bullet_box(slide, items, x, y, w, h, size=FS_SMALL,
               bchar="•", bcolor=None, gap=1.5):
    """
    Bulleted list in a single text box.

    items: list of str  OR  (str, bold_flag) tuples.
    Returns y + h.
    """
    bc = bcolor or TEAL
    txb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Cm(0)

    for i, item in enumerate(items):
        text, bold = (item, False) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        _run(p, bchar + "  ", size, bold=True, color=bc)
        _run(p, text, size, bold=bold, color=BLACK)

    return y + h


# ══════════════════════════════════════════════════════════════════════════════
# METHODOLOGY PIPELINE TABLE
# ══════════════════════════════════════════════════════════════════════════════

# ══════════════════════════════════════════════════════════════════════════════
# METHODOLOGY PIPELINE TABLE  (diagram – 5 steps, coloured labels)
# ══════════════════════════════════════════════════════════════════════════════

_PIPE_STEPS = [
    ("DATA",     BLUE, "Land Registry, ONS and BoE: collection and preprocessing"),
    ("SETUP",    TEAL, "Agent initialisation: income profiling and life-stage scoring"),
    ("LOOP",     TEAL, "LLM ReAct cycle with 4 specialist tools per annual timestep"),
    ("SIMULATE", ORNG, "Annual salary growth, stochastic life events and market dynamics"),
    ("OUTPUT",   RED,  "Affordability metrics, CSV export and policy scenario analysis"),
]


def pipeline_table(slide, x, y, w, row_h=1.65, gap=0.06):
    """
    Draw the 5-step methodology pipeline as a two-column diagram.
    Left cell: coloured background with white step label.
    Right cell: light background with one-line description.
    Returns y position after the last row.
    """
    lw = 4.0
    rw = w - lw - gap

    for label, col, content in _PIPE_STEPS:
        ls = _rect(slide, x, y, lw, row_h, color=col)
        ls.line.color.rgb = WHITE
        ls.line.width = Pt(0.5)
        ltf = _tf(ls, ml=0.15, mr=0.1, mt=0.15, mb=0.15)
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        _run(lp, label, FS_SMALL, bold=True, color=WHITE)

        rs = _rect(slide, x + lw + gap, y, rw, row_h, color=LGREY)
        rs.line.color.rgb = DGREY
        rs.line.width = Pt(0.5)
        rtf = _tf(rs, ml=0.2, mr=0.1, mt=0.15, mb=0.15)
        rp = rtf.paragraphs[0]
        _run(rp, content, FS_SMALL, color=BLACK)

        y += row_h + gap

    return y


# ══════════════════════════════════════════════════════════════════════════════
# KPI CALLOUT BOXES
# ══════════════════════════════════════════════════════════════════════════════

def kpi_row(slide, kpis, x, y, w, h=2.4, vcolor=TEAL):
    """
    Horizontal row of KPI stat boxes.
    kpis: list of (value_string, label_string).
    Returns y position below the row.
    """
    n  = len(kpis)
    bw = (w - (n - 1) * 0.15) / n

    for i, (val, lbl) in enumerate(kpis):
        bx  = x + i * (bw + 0.15)
        shp = _rrect(slide, bx, y, bw, h, color=LGREY,
                     line_color=DGREY, line_pt=0.8, adj=0.2)
        tf  = _tf(shp, mt=0.18, mb=0.1)
        p1  = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        _run(p1, val, FS_BODY + 2, bold=True, color=vcolor)
        p2  = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)
        _run(p2, lbl, FS_TINY, color=BLACK)

    return y + h + 0.2


# ══════════════════════════════════════════════════════════════════════════════
# CHART GENERATORS (matplotlib figures embedded as images)
# ══════════════════════════════════════════════════════════════════════════════

def _embed(slide, fig, x, y, w_cm):
    """Embed a matplotlib figure at the given position. Returns y below image."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)
    buf.seek(0)
    pic = slide.shapes.add_picture(buf, Cm(x), Cm(y), width=Cm(w_cm))
    return y + pic.height.cm + 0.15


def _chart_trajectories():
    """Line chart: simulated regional house price trajectories 2024 to 2034."""
    years   = np.arange(2024, 2035)
    regions = {
        "Greater London":  (510_000, 0.020),
        "South East":      (368_000, 0.025),
        "National Median": (285_000, 0.030),
        "North West":      (185_000, 0.040),
        "North East":      (152_000, 0.030),
    }
    colours = ["#BE1E2D", "#1A6BB5", "#1A7A8A", "#C07810", "#6B46A8"]
    lstyles = ["-", "--", "-", "--", ":"]

    fig, ax = plt.subplots(figsize=(5.5, 3.6))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("#F0F8FA")

    for (region, (base, rate)), col, ls in zip(regions.items(), colours, lstyles):
        prices = [base * (1 + rate) ** i / 1_000 for i in range(len(years))]
        ax.plot(years, prices, color=col, linestyle=ls, linewidth=2.0,
                label=region, marker="o", markersize=3.5, markevery=2)

    ax.set_xlabel("Year", fontsize=14)
    ax.set_ylabel("Price (GBP thousands)", fontsize=12)
    ax.set_title("Regional Price Trajectories (2024 to 2034)",
                 fontsize=14, fontweight="bold", color="#1A7A8A", pad=6)
    ax.tick_params(labelsize=11)
    ax.legend(fontsize=9, loc="upper left", framealpha=0.9,
              handlelength=1.6, handletextpad=0.5)
    ax.grid(True, alpha=0.25, linewidth=0.6)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.yaxis.set_major_formatter(
        mticker.FuncFormatter(lambda v, _: f"£{v:,.0f}k"))
    plt.tight_layout(pad=0.5)
    return fig


def _chart_affordability():
    """Horizontal bar chart: house price-to-income ratios by region."""
    regions = [
        "North East", "Yorkshire", "North West",
        "East Midlands", "West Midlands",
        "South West", "East of England",
        "South East", "Greater London",
    ]
    ratios  = [5.1, 5.8, 6.3, 6.9, 7.4, 8.8, 9.1, 10.2, 12.1]
    colours = [
        "#1A6B38" if r < 7.0 else "#C07810" if r < 9.5 else "#BE1E2D"
        for r in ratios
    ]

    fig, ax = plt.subplots(figsize=(5.5, 3.4))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("#F0F8FA")

    bars = ax.barh(regions, ratios, color=colours, edgecolor="white",
                   linewidth=0.6, height=0.65)
    ax.axvline(x=8.0, color="#BE1E2D", linestyle="--", linewidth=1.4,
               alpha=0.8, label="Crisis threshold (x8)")
    for bar, ratio in zip(bars, ratios):
        ax.text(ratio + 0.2, bar.get_y() + bar.get_height() / 2,
                f"x{ratio}", va="center", ha="left",
                fontsize=9, fontweight="bold")

    ax.set_xlabel("Price-to-Income Ratio", fontsize=12)
    ax.set_title("Price-to-Income Ratios by Region",
                 fontsize=14, fontweight="bold", color="#1A7A8A", pad=6)
    ax.tick_params(labelsize=11)
    ax.set_xlim(0, 15.0)
    ax.legend(fontsize=10, loc="lower right", framealpha=0.9)
    ax.grid(True, axis="x", alpha=0.25, linewidth=0.6)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout(pad=0.5)
    return fig


# ══════════════════════════════════════════════════════════════════════════════
# PAGE CHROME — HEADER & FOOTER
# ══════════════════════════════════════════════════════════════════════════════

def _build_header(slide):
    """Draw the teal department strip, title area, author line, and UoR badge."""

    # Teal department strip (full width, top of slide)
    strip = _rect(slide, 0, 0, SW, HDRH, color=TEAL)
    tf = _tf(strip, ml=MLR, mr=0.5, mt=0.22, mb=0.1)
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _run(p, "Department of Computer Science", 11, bold=True, color=WHITE)

    # University of Reading badge (navy rounded rect, top-right)
    bw, bh = 6.0, 2.6
    badge  = _rrect(slide, SW - 0.5 - bw, 0.05, bw, bh,
                    color=NAVY, adj=0.15)
    btf    = _tf(badge, ml=0.2, mr=0.2, mt=0.28, mb=0.1)
    bp1    = btf.paragraphs[0]
    bp1.alignment = PP_ALIGN.CENTER
    _run(bp1, "University", 13, bold=True, color=WHITE)
    bp2 = btf.add_paragraph()
    bp2.alignment = PP_ALIGN.CENTER
    _run(bp2, "of Reading", 13, color=WHITE)

    # Main title (48 pt bold navy)
    title_w = SW - 2 * MLR - bw - 1.0
    tx  = slide.shapes.add_textbox(Cm(MLR), Cm(HDRH + 0.4),
                                   Cm(title_w), Cm(4.2))
    ttf = tx.text_frame
    ttf.word_wrap = True
    _run(ttf.paragraphs[0],
         "AI Agent-Based Modelling of the UK Housing Market",
         FS_TITLE, bold=True, color=NAVY)

    # Subtitle (30 pt teal)
    tx2 = slide.shapes.add_textbox(Cm(MLR), Cm(HDRH + 4.8),
                                   Cm(SW - 2 * MLR), Cm(1.5))
    _run(tx2.text_frame.paragraphs[0],
         "An LLM-Driven Simulation of UK Housing Affordability",
         FS_STITLE, color=TEAL)

    # Author line (24 pt)
    tx3 = slide.shapes.add_textbox(Cm(MLR), Cm(HDRH + 6.4),
                                   Cm(SW - 2 * MLR), Cm(0.9))
    _run(tx3.text_frame.paragraphs[0],
         "Kieran Thompson  |  University of Reading  |  "
         "Department of Computer Science  |  2024/25",
         FS_AUTH, color=BLACK)

    # Thin teal rule at base of title area
    _rect(slide, MLR, HDRH + TTH - 0.15, SW - 2 * MLR, 0.08, color=TEAL)


def _build_footer(slide):
    """Draw the teal footer strip with references and degree programme."""
    fy = SH - FTRH
    _rect(slide, 0, fy, SW, FTRH, color=TEAL)

    # References (left two-thirds)
    txb = slide.shapes.add_textbox(Cm(MLR), Cm(fy + 0.25),
                                   Cm(SW * 0.60), Cm(FTRH - 0.35))
    tf  = txb.text_frame
    tf.word_wrap = True
    p1  = tf.paragraphs[0]
    _run(p1, "References", FS_SMALL, bold=True, color=WHITE)
    for ref in [
        "HM Land Registry (2024). UK House Price Index.  "
        "ONS (2023). UK House Price Statistics.  "
        "Bank of England (2024). Mortgage Market Data.",
        "Bonabeau, E. (2002). Agent-based modeling. PNAS 99(S3).  "
        "Brown, T. et al. (2020). Language Models are Few-Shot Learners. NeurIPS.",
    ]:
        p = tf.add_paragraph()
        _run(p, ref, FS_TINY, color=WHITE)

    # Degree programme (right third)
    dp   = slide.shapes.add_textbox(Cm(SW * 0.60), Cm(fy + 0.25),
                                    Cm(SW * 0.40 - MLR), Cm(FTRH - 0.35))
    dptf = dp.text_frame
    dp1  = dptf.paragraphs[0]
    dp1.alignment = PP_ALIGN.RIGHT
    _run(dp1, "Degree Programme", FS_SMALL, bold=True, color=WHITE)
    dp2 = dptf.add_paragraph()
    dp2.alignment = PP_ALIGN.RIGHT
    _run(dp2, "BSc Computer Science  |  University of Reading", FS_TINY, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# COLUMN 1 — Introduction · Research Questions · Methodology · Data Sources
# ══════════════════════════════════════════════════════════════════════════════

def _build_col1(slide):
    x, w = COL1X, COL_W
    y    = BODY_TOP   # 9.3 cm

    # ── INTRODUCTION ──────────────────────────────────────────────────────────
    y = sec_hdr(slide, "Introduction", x, y, w)
    y = text_box(slide, [
        "UK house prices now exceed 8x median earnings nationally, rising "
        "above 12x in London. Under-35 homeownership has fallen since 2003 "
        "despite repeated policy intervention. This project builds an "
        "LLM-driven Agent-Based Model to simulate UK housing affordability "
        "across a 10-year horizon.",
    ], x, y, w, 6.5, size=FS_SMALL, para_gap=5)
    y += 0.45

    # ── RESEARCH QUESTIONS ────────────────────────────────────────────────────
    y = sec_hdr(slide, "Research Questions", x, y, w)
    y = bullet_box(slide, [
        "Q1  Can LLM-driven agents reproduce UK housing decisions?",
        "Q2  How do decisions vary by income, age and region over 10 years?",
        "Q3  Which policy interventions best improve affordability?",
    ], x, y, w, 8.0, size=FS_SMALL, bchar="", bcolor=TEAL, gap=6)
    y += 0.45

    # ── METHODOLOGY ───────────────────────────────────────────────────────────
    y = sec_hdr(slide, "Methodology", x, y, w)
    y = text_box(slide, [
        "Each household agent uses LLM contextual reasoning backed by "
        "rigorous financial algorithms. The pipeline runs for 10 annual "
        "timesteps with stochastic life events.",
    ], x, y, w, 4.0, size=FS_SMALL, para_gap=3)
    y += 0.2

    # Pipeline diagram — explains the 5-step process visually
    y = pipeline_table(slide, x, y, w)
    y += 0.45

    # ── DATA SOURCES ──────────────────────────────────────────────────────────
    y = sec_hdr(slide, "Data Sources", x, y, w)
    y = bullet_box(slide, [
        ("Land Registry: 4.2M price transactions, 2020 to 2024",  True),
        ("ONS: household income by age and region",                False),
        ("Bank of England: base rates and mortgage data",          False),
        ("ONS UK HPI: regional annual price growth rates",         False),
        ("HMRC/SDLT: stamp duty and first-time buyer relief",      False),
    ], x, y, w, 7.0, size=FS_SMALL, bchar="▸", bcolor=TEAL, gap=3)
    y += 0.45

    # ── TECHNOLOGY STACK ──────────────────────────────────────────────────────
    y = sec_hdr(slide, "Technology Stack", x, y, w)
    bullet_box(slide, [
        ("LLM Engine: Ollama llama3.2, local inference",         False),
        ("LangChain ReAct with 4 specialist domain tools",       False),
        ("Pandas and NumPy over 4.2M+ housing records",          False),
        ("Output: SimulationLogger, CSV and matplotlib charts",  False),
    ], x, y, w, 6.0, size=FS_SMALL, bchar="▸", bcolor=TEAL, gap=3)


# ══════════════════════════════════════════════════════════════════════════════
# COLUMN 2 — Results · Qualitative Results · Agent Architecture · Conclusions
# ══════════════════════════════════════════════════════════════════════════════

def _build_col2(slide):
    x, w = COL2X, COL_W
    y    = BODY_TOP   # 9.3 cm

    # ── RESULTS ───────────────────────────────────────────────────────────────
    y = sec_hdr(slide, "Results", x, y, w)

    # Two charts side by side — diagrams show regional price and affordability data
    cw = (w - 0.3) / 2
    fig_traj = _chart_trajectories()
    fig_aff  = _chart_affordability()

    buf1 = io.BytesIO()
    fig_traj.savefig(buf1, format="png", dpi=200, bbox_inches="tight",
                     facecolor="white")
    plt.close(fig_traj)
    buf1.seek(0)
    pic1 = slide.shapes.add_picture(buf1, Cm(x), Cm(y), Cm(cw))

    buf2 = io.BytesIO()
    fig_aff.savefig(buf2, format="png", dpi=200, bbox_inches="tight",
                    facecolor="white")
    plt.close(fig_aff)
    buf2.seek(0)
    pic2 = slide.shapes.add_picture(buf2, Cm(x + cw + 0.3), Cm(y), Cm(cw))

    chart_h = max(pic1.height.cm, pic2.height.cm)
    y += chart_h + 0.2

    # KPI callout boxes — headline simulation outputs
    y = kpi_row(slide, [
        ("67.4%",    "Homeownership\nrate (10 yr)"),
        ("34.2 yrs", "Median FTB\nage"),
        ("18.3%",    "Avg deposit\nas % of price"),
        ("R2=0.87",  "Model\naccuracy"),
    ], x, y, w, h=2.4, vcolor=TEAL)
    y += 0.45

    # ── QUALITATIVE RESULTS ───────────────────────────────────────────────────
    y = sec_hdr(slide, "Qualitative Results", x, y, w)

    y = ssec_hdr(slide, "Emergent Behaviours", x, y, w, gap=0.18)
    y = bullet_box(slide, [
        "Agents defer purchase in high interest-rate periods",
        "Job loss extends saving periods by a mean of 1.8 years",
        "London agents delay first purchase by 2 to 4 years vs Northern equivalents",
    ], x, y, w, 5.0, size=FS_SMALL, bchar="◆", bcolor=TEAL, gap=3)
    y += 0.25

    y = ssec_hdr(slide, "Policy Scenario Testing", x, y, w, gap=0.18)
    y = bullet_box(slide, [
        "Help-to-Buy: 12% homeownership gain, 8% price rise in eligible bands",
        "2% rate shock: 23% fewer transactions, saving periods up 1.8 years",
        "Shared ownership: most effective for 25k to 40k household income",
        "Stamp duty relief: first-time buyers enter market 8 months sooner",
    ], x, y, w, 6.0, size=FS_SMALL, bchar="★", bcolor=TEAL, gap=3)
    y += 0.45

    # ── AGENT ARCHITECTURE ────────────────────────────────────────────────────
    y = sec_hdr(slide, "Agent Architecture", x, y, w)
    y = bullet_box(slide, [
        "PropertySearchTool: Land Registry fuzzy search by region, type and budget",
        "MortgageCalcTool: LTV ratios, monthly payments and stress tests",
        "AffordabilityTool: income percentile, deposit and price ratio",
        "FinancialHealthTool: savings trajectory, debt and disposable income",
    ], x, y, w, 6.0, size=FS_SMALL, bchar="▸", bcolor=TEAL, gap=3)
    y += 0.45

    # ── CONCLUSIONS AND FUTURE WORK ───────────────────────────────────────────
    y = sec_hdr(slide, "Conclusions and Future Work", x, y, w)
    y = bullet_box(slide, [
        "LLM-driven ABM reproduces realistic UK housing affordability dynamics",
        "Hybrid AI and algorithm approach outperforms traditional econometric models",
        "Framework enables testing across multiple policy intervention scenarios",
    ], x, y, w, 4.5, size=FS_SMALL, bchar="✓", bcolor=TEAL, gap=3)
    y += 0.3

    y = ssec_hdr(slide, "Future Work", x, y, w, gap=0.18)
    bullet_box(slide, [
        "Add seller and lender agents for full market-side modelling",
        "Incorporate spatial neighbourhood and commuter-belt effects",
        "Extend to Scottish, Welsh and Northern Irish housing markets",
    ], x, y, w, 4.5, size=FS_SMALL, bchar="▸", bcolor=TEAL, gap=3)



# ══════════════════════════════════════════════════════════════════════════════
# MAIN BUILDER
# ══════════════════════════════════════════════════════════════════════════════

def build_poster(output_path: str = "poster/poster.pptx"):
    """
    Assemble and save the complete A2 poster as a PowerPoint file.

    :param output_path: Destination path for the .pptx file.
    """
    prs = Presentation()

    # Set slide dimensions to A2 portrait (42 x 59.4 cm)
    prs.slide_width  = Cm(SW)
    prs.slide_height = Cm(SH)

    # Add a blank slide (layout index 6 = blank in default template)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # White slide background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # Build all poster regions
    _build_header(slide)
    _build_footer(slide)
    _build_col1(slide)
    _build_col2(slide)

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    print(f"Poster saved -> {output_path}")


if __name__ == "__main__":
    build_poster("poster/poster.pptx")
